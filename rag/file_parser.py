"""
文件解析模块 - 用于 Moodle 附件分析

使用轻量级库支持多种文档格式，提取文本内容供 LLM 分析

依赖:
    pip install PyPDF2 python-docx python-pptx markdown
"""

import logging
from pathlib import Path
from typing import Optional, Dict, Any, List
import mimetypes

logger = logging.getLogger(__name__)


class FileParser:
    """文件解析器 - 使用轻量级库支持多种文档格式"""
    
    SUPPORTED_EXTENSIONS = {
        '.pdf', '.ppt', '.pptx', '.docx', '.md', '.txt'
    }
    
    def __init__(self):
        """初始化文件解析器"""
        pass  # 不需要预先检查依赖，按需导入
    
    def parse_file(
        self,
        file_path: str,
        max_chars: Optional[int] = None,
        extract_tables: bool = True,
        extract_images: bool = False
    ) -> Dict[str, Any]:
        """
        解析文件并提取文本内容
        
        Args:
            file_path: 文件路径
            max_chars: 最大字符数限制（用于 LLM token 控制）
            extract_tables: 是否提取表格内容
            extract_images: 是否提取图片信息
        
        Returns:
            解析结果字典，包含:
            - success: 是否成功
            - text: 提取的纯文本
            - elements: 结构化元素列表
            - metadata: 文件元数据
            - error: 错误信息（如果失败）
        """
        path = Path(file_path)
        
        if not path.exists():
            return {
                "success": False,
                "error": f"文件不存在: {file_path}",
                "text": "",
                "elements": []
            }
        
        file_ext = path.suffix.lower()
        
        if file_ext not in self.SUPPORTED_EXTENSIONS:
            return {
                "success": False,
                "error": f"不支持的文件格式: {file_ext}",
                "text": "",
                "elements": [],
                "supported_formats": list(self.SUPPORTED_EXTENSIONS)
            }
        
        try:
            # 根据文件类型选择解析策略
            if file_ext == '.pdf':
                full_text = self._parse_pdf(path, extract_tables)
            elif file_ext in ['.ppt', '.pptx']:
                full_text = self._parse_ppt(path)
            elif file_ext in ['.doc', '.docx']:
                full_text = self._parse_docx(path)
            elif file_ext == '.md':
                full_text = self._parse_markdown(path)
            elif file_ext == '.txt':
                full_text = self._parse_text(path)
            else:
                raise ValueError(f"不支持的文件格式: {file_ext}")
            
            original_length = len(full_text)
            
            # 限制字符数
            if max_chars and len(full_text) > max_chars:
                full_text = full_text[:max_chars] + f"\n\n... (文本已截断，原始长度: {original_length} 字符)"
            
            # 序列化元素
            elements_list = self._serialize_elements(full_text)
            
            # 构建结果
            result = {
                "success": True,
                "text": full_text,
                "char_count": len(full_text),
                "original_char_count": original_length,
                "truncated": original_length > (max_chars or float('inf')),
                "element_count": len(elements_list),
                "file_name": path.name,
                "file_size_kb": path.stat().st_size / 1024,
                "file_type": file_ext,
                "elements": elements_list
            }
            
            logger.info(f"✅ 成功解析 {path.name}: {original_length} 字符")
            
            return result
            
        except Exception as e:
            logger.error(f"❌ 解析文件失败 {path.name}: {e}", exc_info=True)
            return {
                "success": False,
                "error": str(e),
                "text": "",
                "elements": []
            }
    
    def _parse_pdf(self, path: Path, extract_tables: bool = True) -> str:
        """解析 PDF 文件 - 使用 PyPDF2"""
        try:
            from PyPDF2 import PdfReader
        except ImportError:
            raise ImportError("请安装 PyPDF2: pip install PyPDF2")
        
        logger.info(f"📄 正在解析 PDF: {path.name}")
        
        reader = PdfReader(str(path))
        texts = []
        
        for page_num, page in enumerate(reader.pages, 1):
            try:
                text = page.extract_text()
                if text.strip():
                    texts.append(f"=== 第 {page_num} 页 ===\n{text}")
            except Exception as e:
                logger.warning(f"⚠️  页面 {page_num} 提取失败: {e}")
        
        return "\n\n".join(texts)
    
    def _parse_ppt(self, path: Path) -> str:
        """解析 PPT/PPTX 文件 - 使用 python-pptx"""
        try:
            from pptx import Presentation
        except ImportError:
            raise ImportError("请安装 python-pptx: pip install python-pptx")
        
        logger.info(f"📊 正在解析 PowerPoint: {path.name}")
        
        prs = Presentation(str(path))
        texts = []
        
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_texts = [f"=== 幻灯片 {slide_num} ==="]
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    slide_texts.append(shape.text)
            
            if len(slide_texts) > 1:
                texts.append("\n".join(slide_texts))
        
        return "\n\n".join(texts)
    
    def _parse_docx(self, path: Path) -> str:
        """解析 Word 文件 - 使用 python-docx"""
        try:
            from docx import Document
        except ImportError:
            raise ImportError("请安装 python-docx: pip install python-docx")
        
        logger.info(f"📝 正在解析 Word 文档: {path.name}")
        
        doc = Document(str(path))
        texts = []
        
        # 提取段落
        for para in doc.paragraphs:
            if para.text.strip():
                texts.append(para.text)
        
        # 提取表格
        for table in doc.tables:
            table_texts = []
            for row in table.rows:
                row_text = " | ".join(cell.text.strip() for cell in row.cells)
                if row_text.strip():
                    table_texts.append(row_text)
            if table_texts:
                texts.append("\n表格内容:\n" + "\n".join(table_texts))
        
        return "\n\n".join(texts)
    
    def _parse_markdown(self, path: Path) -> str:
        """解析 Markdown 文件"""
        logger.info(f"📖 正在解析 Markdown: {path.name}")
        
        # Markdown 直接读取原始文本即可
        return path.read_text(encoding='utf-8')
    
    def _parse_text(self, path: Path) -> str:
        """解析纯文本文件"""
        logger.info(f"📃 正在解析文本文件: {path.name}")
        
        # 尝试多种编码
        encodings = ['utf-8', 'gbk', 'gb2312', 'latin-1']
        
        for encoding in encodings:
            try:
                return path.read_text(encoding=encoding)
            except UnicodeDecodeError:
                continue
        
        # 如果都失败，使用 errors='ignore'
        return path.read_text(encoding='utf-8', errors='ignore')
    
    def _extract_text_from_elements(self, text: str) -> str:
        """兼容接口 - 直接返回文本"""
        return text
    
    def _serialize_elements(self, text: str) -> List[Dict[str, Any]]:
        """序列化为元素列表（简化版）"""
        # 简单按段落分割
        paragraphs = [p.strip() for p in text.split('\n\n') if p.strip()]
        
        return [
            {
                "type": "Paragraph",
                "text": para,
                "category": "text"
            }
            for para in paragraphs
        ]
    
    def get_file_info(self, file_path: str) -> Dict[str, Any]:
        """获取文件基本信息"""
        path = Path(file_path)
        
        if not path.exists():
            return {
                "exists": False,
                "error": "文件不存在"
            }
        
        file_ext = path.suffix.lower()
        mime_type, _ = mimetypes.guess_type(str(path))
        
        return {
            "exists": True,
            "name": path.name,
            "extension": file_ext,
            "mime_type": mime_type,
            "size_bytes": path.stat().st_size,
            "size_kb": path.stat().st_size / 1024,
            "size_mb": path.stat().st_size / (1024 * 1024),
            "supported": file_ext in self.SUPPORTED_EXTENSIONS
        }


def parse_file_for_llm(
    file_path: str,
    max_chars: int = 30000,
    extract_tables: bool = True
) -> str:
    """
    便捷函数：解析文件并返回适合 LLM 的文本
    
    Args:
        file_path: 文件路径
        max_chars: 最大字符数（考虑 LLM context 限制）
        extract_tables: 是否提取表格
    
    Returns:
        提取的文本内容
    
    Raises:
        Exception: 解析失败时抛出异常
    """
    parser = FileParser()
    result = parser.parse_file(
        file_path=file_path,
        max_chars=max_chars,
        extract_tables=extract_tables
    )
    
    if not result["success"]:
        raise Exception(f"文件解析失败: {result.get('error', 'Unknown error')}")
    
    return result["text"]


def analyze_file_with_llm(
    file_path: str,
    analysis_prompt: Optional[str] = None,
    max_file_chars: int = 20000
) -> Dict[str, Any]:
    """
    解析文件并使用 LLM 分析内容
    
    Args:
        file_path: 文件路径
        analysis_prompt: 自定义分析提示（可选）
        max_file_chars: 文件内容最大字符数
    
    Returns:
        LLM 分析结果
    """
    from openai import OpenAI
    import os
    
    # 检查 API Key
    api_key = os.getenv("OPENAI_API_KEY")
    if not api_key:
        return {
            "success": False,
            "error": "OPENAI_API_KEY 环境变量未设置"
        }
    
    # 解析文件
    parser = FileParser()
    parse_result = parser.parse_file(
        file_path=file_path,
        max_chars=max_file_chars,
        extract_tables=True
    )
    
    if not parse_result["success"]:
        return {
            "success": False,
            "error": f"文件解析失败: {parse_result.get('error')}"
        }
    
    file_text = parse_result["text"]
    file_name = parse_result["file_name"]
    
    # 默认分析提示
    if not analysis_prompt:
        analysis_prompt = f"""
请分析以下文档内容，提取关键信息：

文件名: {file_name}

文档内容:
{file_text}

请提取：
1. 主要内容摘要（3-5句话）
2. 关键要点和重要信息
3. 如果是课程作业文档，请提取：
   - 截止日期
   - 作业要求
   - 评分标准
   - 注意事项
4. 其他重要信息

以结构化的方式返回。
"""
    else:
        # 使用自定义提示，但确保包含文件内容
        analysis_prompt = f"""
文件名: {file_name}

文档内容:
{file_text}

{analysis_prompt}
"""
    
    try:
        # 调用 LLM
        base_url = os.getenv("OPENAI_BASE_URL")
        model = os.getenv("OPENAI_MODEL", "gpt-4o-mini")
        
        if base_url:
            client = OpenAI(api_key=api_key, base_url=base_url)
        else:
            client = OpenAI(api_key=api_key)
        
        logger.info(f"🤖 正在使用 {model} 分析文档...")
        
        response = client.chat.completions.create(
            model=model,
            messages=[
                {
                    "role": "system",
                    "content": "你是一个智能文档分析助手，擅长从各类文档中提取关键信息和结构化数据。"
                },
                {
                    "role": "user",
                    "content": analysis_prompt
                }
            ],
            temperature=0.3,
            max_tokens=2000
        )
        
        analysis_result = response.choices[0].message.content
        
        return {
            "success": True,
            "file_name": file_name,
            "file_type": parse_result["file_type"],
            "char_count": parse_result["char_count"],
            "element_count": parse_result["element_count"],
            "analysis": analysis_result,
            "model": model,
            "tokens_used": {
                "prompt": response.usage.prompt_tokens,
                "completion": response.usage.completion_tokens,
                "total": response.usage.total_tokens
            }
        }
        
    except Exception as e:
        logger.error(f"❌ LLM 分析失败: {e}", exc_info=True)
        return {
            "success": False,
            "error": f"LLM 分析失败: {str(e)}",
            "file_text": file_text  # 返回文本供后续处理
        }


# 便捷函数集合

def extract_text_from_pdf(pdf_path: str, max_chars: Optional[int] = None) -> str:
    """从 PDF 提取文本"""
    return parse_file_for_llm(pdf_path, max_chars=max_chars or 50000)


def extract_text_from_docx(docx_path: str, max_chars: Optional[int] = None) -> str:
    """从 Word 文档提取文本"""
    return parse_file_for_llm(docx_path, max_chars=max_chars or 50000)


def extract_text_from_pptx(pptx_path: str, max_chars: Optional[int] = None) -> str:
    """从 PowerPoint 提取文本"""
    return parse_file_for_llm(pptx_path, max_chars=max_chars or 50000)


def batch_parse_files(file_paths: List[str]) -> List[Dict[str, Any]]:
    """批量解析多个文件"""
    parser = FileParser()
    results = []
    
    for file_path in file_paths:
        result = parser.parse_file(file_path)
        results.append({
            "file": file_path,
            **result
        })
    
    return results


if __name__ == "__main__":
    """测试代码"""
    import sys
    
    if len(sys.argv) < 2:
        print("用法: python file_parser.py <file_path> [--llm]")
        print("示例: python file_parser.py test.pdf")
        print("示例: python file_parser.py test.docx --llm")
        sys.exit(1)
    
    file_path = sys.argv[1]
    enable_llm = "--llm" in sys.argv
    
    parser = FileParser()
    
    # 获取文件信息
    info = parser.get_file_info(file_path)
    print(f"\n📁 文件信息:")
    for key, value in info.items():
        print(f"  {key}: {value}")
    
    # 解析文件
    result = parser.parse_file(file_path, max_chars=30000)
    
    if result["success"]:
        print(f"\n✅ 解析成功！")
        print(f"  - 元素数量: {result['element_count']}")
        print(f"  - 字符数量: {result['char_count']}")
        print(f"\n📄 文本预览 (前500字符):")
        print("-" * 60)
        print(result["text"][:500])
        print("-" * 60)
        
        # LLM 分析
        if enable_llm:
            print("\n🤖 正在调用 LLM 分析文档...")
            llm_result = analyze_file_with_llm(file_path)
            
            if llm_result["success"]:
                print(f"\n✅ LLM 分析完成 (用了 {llm_result['tokens_used']['total']} tokens)")
                print("\n📊 分析结果:")
                print("-" * 60)
                print(llm_result["analysis"])
                print("-" * 60)
            else:
                print(f"\n❌ LLM 分析失败: {llm_result['error']}")
    else:
        print(f"\n❌ 解析失败: {result['error']}")

